
import os
import io
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
from PIL import Image
try:
    import pytesseract
except ImportError:
    pytesseract = None

# Note: pytesseract requires Tesseract-OCR installed on the system.
# For simplicity, if not installed, we might skip OCR or just use basic text extraction.
# Or better, we can send Image bytes directly to Gemini Vision!
# So for images, we won't extract text here, we'll return the image object/bytes.

def extract_text(file_storage):
    """
    Extracts text from the uploaded file (Werkzeug FileStorage).
    Returns a string of text, or a list of images/blobs if it's an image.
    """
    filename = file_storage.filename.lower()
    
    try:
        if filename.endswith('.pdf'):
            # Return properly tailored object for both Gemini (bytes) and Offline (text)
            try:
                # 1. Get Text for Fallback
                reader = PdfReader(file_storage.stream)
                text = ""
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
                
                # Reset stream to read bytes
                file_storage.seek(0)
                
                return {
                    "type": "pdf", 
                    "content": file_storage.read(), 
                    "mime_type": "application/pdf",
                    "fallback_text": text
                }
            except Exception as e:
                print(f"PDF Parse Error: {e}")
                return {"type": "error", "content": "Failed to parse PDF"}

        elif filename.endswith('.docx'):
            doc = Document(file_storage.stream)
            text = "\n".join([para.text for para in doc.paragraphs])
            return {"type": "text", "content": text}

        elif filename.endswith('.pptx'):
            prs = Presentation(file_storage.stream)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return {"type": "text", "content": text}

        elif filename.endswith('.xlsx'):
            wb = openpyxl.load_workbook(file_storage.stream, data_only=True)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += " ".join([str(c) for c in row if c is not None]) + "\n"
            return {"type": "text", "content": text}
            
        elif filename.endswith(('.png', '.jpg', '.jpeg')):
            # For images, we return the stream content to send to Gemini Vision
            # We reset the stream pointer just in case
            file_storage.stream.seek(0)
            return {"type": "image", "content": file_storage.read(), "mime_type": file_storage.mimetype}
            
        elif filename.endswith('.txt'):
            return {"type": "text", "content": file_storage.read().decode('utf-8', errors='ignore')}

        else:
            return {"type": "error", "content": "Unsupported file format."}
            
    except Exception as e:
        return {"type": "error", "content": str(e)}
